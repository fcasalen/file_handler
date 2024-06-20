from pptx import Presentation
from .utils import error_loading
from .validators import TxtData

def extract_text_from_ppt(ppt_file):
    presentation = Presentation(ppt_file)
    for slide in presentation.slides:
        text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text.append(run.text)
    return '\n'.join(text)

class PPTHandler:
    @staticmethod
    def load(file_path_with_password:tuple[str, str], encoding:str, mode:str = 'r'):
        file_path, password = file_path_with_password
        slides = {}
        try:
            presentation = Presentation(file_path)
            text = []
            for numSlide, slide in enumerate(presentation.slides, start=1):
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text.append(run.text)
                slides[f'Slide {numSlide}'] = '\n'.join(text)
            return slides
        except Exception as error:
            return error_loading(file_path, error=error)

    @staticmethod
    def write(file_path:str, encoding:str, data:dict[str, str], mode:str = 'w'):
        TxtData(data=data)
        print("I can't write ppt files yet!")