from pathlib import Path

from pptx import Presentation


def load(
    file_path: Path, password: str = None, encoding: str = "utf-8", mode: str = "r"
) -> dict[str, str]:
    """Function to load pptx and ppt files.

    Args:
        file_path (Path): path to the pptx file.
        password (str, optional): password to open the pptx file (not implemented).
            Defaults to None.
        encoding (str, optional): encoding type (not implemented). Defaults to "utf-8".
        mode (str, optional): file mode (not implemented). Defaults to 'r'.

    :returns: dictionary with slide numbers as keys and slide text as values.
    :rtype: dict[str, str]
    """

    slides = {}
    presentation = Presentation(file_path)
    text = []
    for num_slide, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text.append(run.text)
        slides[f"Slide {num_slide}"] = "\n".join(text)
    return slides


def write(
    file_path: str,
    data: dict[str, str],
    password: str = None,
    encoding: str = "utf-8",
    mode: str = "w",
) -> None:
    """Function to write pptx files.

    Args:
        file_path (str): path to save the pptx file.
        data (dict[str, str]): data to be written in the pptx file.
        password (str, optional): password to encrypt the pptx file. Defaults to None.
        encoding (str, optional): encoding type. Defaults to "utf-8".
        mode (str, optional): file mode. Defaults to 'w'.

    Raises:
        NotImplementedError: function not implemented yet.
    """
    raise NotImplementedError("I can't write ppt files yet!")
